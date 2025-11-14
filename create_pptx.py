from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# 創建簡報
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(title, subtitle=""):
    """添加標題頁"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(30, 58, 138)  # 深藍
    
    # 標題
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 副標題
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        p = subtitle_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = RGBColor(200, 200, 200)
        p.alignment = PP_ALIGN.CENTER

def add_content_slide(title, content_list):
    """添加內容頁"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # 標題
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(30, 58, 138)
    
    # 內容
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(5.7))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.space_before = Pt(6)
        p.space_after = Pt(6)

# 第1頁：標題
add_title_slide("教育部SEL Prompt生成系統", "K-12完整版 (國中版完成)")

# 第2頁：項目概述
add_content_slide("📋 項目概述", [
    "🎯 為台灣K-12教師快速生成AI教學Prompt",
    "✅ 已完成：國中版（7-9年級）",
    "🚧 待建：小學版（1-6年級分低中高）",
    "🚧 待建：幼兒園版（3-6歲）",
    "💰 完全免費 | 🌐 GitHub Pages永久部署 | 🎨 支援WIX網站"
])

# 第3頁：國中版完成
add_content_slide("✅ 國中版（已完成）", [
    "5個互動HTML工具",
    "• SEL官方版 | SDG完整版 | 整合版 | 深度版 | 導航主頁",
    "",
    "9份完整文檔",
    "• 部署指南 | WIX整合 | 工具使用 | 政策說明 | Prompt模版",
    "",
    "部署方式",
    "• GitHub Pages（推薦） | WIX整合 | 本地使用"
])

# 第4頁：核心特色
add_content_slide("🌟 核心特色", [
    "✅ 30秒生成教案Prompt",
    "✅ 無需複雜設定，填表單即可",
    "✅ 支援任何LLM（ChatGPT、Claude等）",
    "✅ 下載為.txt檔案",
    "✅ 完全離線可用",
    "✅ 參考教育部相關政策"
])

# 第5頁：小學版規劃
add_content_slide("📚 小學版規劃（待建）", [
    "三個年段各3個工具 + 5份文檔",
    "",
    "低年級（1-3年級）",
    "• 簡化、遊戲化、生活情境",
    "",
    "中年級（3-4年級）",
    "• 中等複雜、社交焦點、綜合活動",
    "",
    "高年級（5-6年級）",
    "• 進階、思考、議題導入"
])

# 第6頁：幼兒園版規劃
add_content_slide("🎈 幼兒園版規劃（待建）", [
    "兩個年段各2-3個工具 + 4份文檔",
    "",
    "3-4歲",
    "• 情緒認識、自我控制、遊戲化",
    "",
    "5-6歲",
    "• 情緒管理、友誼建立、多元理解",
    "",
    "特色：簡化介面、大字體、圖像化"
])

# 第7頁：技術架構
add_content_slide("🏗️ 技術架構", [
    "GitHub Repository結構",
    "├─ elementary/ (小學版 - 9個工具)",
    "├─ kindergarten/ (幼兒園版 - 4個工具)",
    "├─ junior-high/ (國中版 - 5個工具✅)",
    "└─ index.html (K-12導航主頁)",
    "",
    "每個版本：工具 + 文檔 + 使用指南"
])

# 第8頁：開發進度
add_content_slide("📊 開發進度", [
    "✅ 國中版：100%完成",
    "   • 5個工具 | 9份文檔 | 已GitHub部署",
    "",
    "🚧 小學版：0% (待建)",
    "   • 需要：9個工具 | 15份文檔",
    "",
    "🚧 幼兒園版：0% (待建)",
    "   • 需要：4個工具 | 8份文檔"
])

# 第9頁：下一步
add_content_slide("🚀 下一步行動", [
    "1️⃣ 確認小學版設計細節",
    "   • SEL能力簡化方案",
    "   • 教育部課程對應",
    "",
    "2️⃣ 開發小學版",
    "   • 低中高年級各3工具",
    "",
    "3️⃣ 開發幼兒園版",
    "   • 兩個年段各2-3工具"
])

# 第10頁：已有資源
add_content_slide("💾 已有資源（可重用）", [
    "✅ HTML表單結構 & JS邏輯",
    "✅ CSS樣式框架",
    "✅ Prompt生成演算法",
    "✅ Markdown文檔模板",
    "✅ 導航菜單系統",
    "",
    "🎨 可自訂色彩、字體、風格"
])

# 第11頁：成功標誌
add_content_slide("🎯 完成時會有", [
    "✅ 15+ 互動HTML工具",
    "✅ 30+ 完整文檔",
    "✅ 完整K-12教育體系覆蓋",
    "✅ 單一GitHub連結可用",
    "✅ 可整合WIX老師網站",
    "✅ 所有老師直接可用"
])

# 第12頁：備註
add_content_slide("⚠️ 重要聲明", [
    "本工具由教師獨立開發",
    "參考但不代表教育部政策",
    "若有政策更新，請查閱教育部官方資訊",
    "",
    "教育部社會情緒學習相關資訊：",
    "https://www.moe.gov.tw/"
])

# 儲存
prs.save('/mnt/user-data/outputs/SEL_Prompt_Generator_K12_Project.pptx')
print("✅ PPTX已生成：SEL_Prompt_Generator_K12_Project.pptx")
